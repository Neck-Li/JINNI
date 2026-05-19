from __future__ import annotations

from doc2md.converters import register
from doc2md.models.document import Block, BlockType, Document


@register(".pptx")
def convert_pptx(path: str) -> Document:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for .pptx files: pip install python-pptx")

    prs = Presentation(path)
    doc = Document(metadata={"source": path, "slides": len(prs.slides)})

    for slide_num, slide in enumerate(prs.slides, 1):
        # slide title
        if slide.shapes.title and slide.shapes.title.text.strip():
            doc.blocks.append(
                Block(type=BlockType.HEADING, level=2, content=slide.shapes.title.text.strip())
            )

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                md_table = _table_to_md(table)
                doc.blocks.append(Block(type=BlockType.TABLE, content=md_table))
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    if para.level > 0:
                        # indented → list item
                        doc.blocks.append(
                            Block(type=BlockType.LIST_ITEM, level=para.level + 1, content=text)
                        )
                    else:
                        doc.blocks.append(
                            Block(type=BlockType.PARAGRAPH, content=text)
                        )

    return doc


def _table_to_md(table) -> str:
    rows_data = []
    col_count = len(table.columns)

    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        padded = cells + [""] * (col_count - len(cells))
        rows_data.append(padded[:col_count])

    if not rows_data:
        return ""

    lines = ["|" + "|".join(rows_data[0]) + "|"]
    lines.append("|" + "|".join(["---"] * col_count) + "|")
    for row in rows_data[1:]:
        lines.append("|" + "|".join(row) + "|")

    return "\n".join(lines)
